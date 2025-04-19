"""
Core layout engine – per‑child units, one shared add_box().
"""
from __future__ import annotations
from dataclasses import dataclass
from typing import Generator, List, Optional, Literal, Sequence

from loguru import logger
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.presentation import Presentation
from pptx.util import Emu, Pt
import polars as pl

from .errors import LayoutError, LayoutStateError, PPTXError, SpecMismatchError
from .units import Size, Auto, resolve_length_span
from .content import Table


@dataclass(slots=True)
class Rect:
    """A rectangle in EMUs with position and size."""
    x: Emu
    y: Emu
    width: Emu
    height: Emu

    def split_horizontal(self, lengths: List[Size]) -> list[Rect]:
        widths = resolve_length_span(lengths, self.width)
        out: list[Rect] = []
        cursor = self.x
        for w in widths:
            out.append(Rect(Emu(cursor), self.y, w, self.height))
            cursor += w
        return out

    def split_vertical(self, lengths: List[Size]) -> List[Rect]:
        heights = resolve_length_span(lengths, self.height)
        out: list[Rect] = []
        cursor = self.y
        for h in heights:
            out.append(Rect(self.x, Emu(cursor), self.width, h))
            cursor += h
        return out

class Area:
    """A layout box that can be split vertically or horizontally."""

    def __init__(self, length: Optional[Size] = None) -> None:
        self.children: List[Area] = []
        self._rect: Optional[Rect] = None
        self.length: Optional[Size] = length
        self._parent: Optional[Area] = None
        self.split_direction: Optional[Literal["vertical", "horizontal", "v", "h"]] = None

    @property
    def rect(self) -> Rect:
        if self._rect is None:
            raise LayoutStateError("Box has no rect; call resolve() first")
        return self._rect

    @property
    def parent(self) -> Area:
        if self._parent is None:
            raise LayoutStateError("Box has no parent")
        return self._parent

    @property
    def root(self) -> SlideRoot:
        *_, last = self.walk_up()
        if isinstance(last, SlideRoot):
            return last
        raise LayoutStateError("No SlideLayout root found in hierarchy")

    @property
    def parent_pos(self) -> int:
        """Get index of current area in siblings."""
        return self.parent.children.index(self)

    @property
    def pos(self) -> tuple[int,...]:
        """Get list of indicies from slideroot."""
        idx: list[int] = []
        for level in self.walk_up():
            if isinstance(level, SlideRoot):
                break
            idx.append(level.parent_pos)
        return tuple(idx[::-1])

    def _get_child(self, index: int) -> Area:
        if index < 0 or index >= len(self.children):
            raise IndexError("Index out of range")
        return self.children[index]

    def __getitem__(self, index: int|Sequence[int]) -> Area:
        if isinstance(index, int):
            return self._get_child(index)
        elif len(index) == 0:
            return self
        else:
            return self._get_child(index[0])[index[1:]]

    def add_child(self, item: Area) -> Area:
        item._parent = self
        self.children.append(item)
        return item

    def walk(self) -> Generator[Area]:
        yield self
        for child in self.children:
            yield from child.walk()

    def walk_up(self) -> Generator[Area]:
        yield self
        if isinstance(self, SlideRoot):
            return
        yield from self.parent.walk_up()

    def layout(self, rect: Rect):
        logger.debug(f"{len(list(self.walk_up()))} - {rect}")
        self._rect = rect
        self._layout_children()

    def _layout_children(self):
        if self.split_direction == "vertical":
            lengths = [child.length or Auto(1) for child in self.children]
            rects = self.rect.split_vertical(lengths)
        elif self.split_direction == "horizontal":
            lengths = [child.length or Auto(1) for child in self.children]
            rects = self.rect.split_horizontal(lengths)
        else:
            return
        for child, rect in zip(self.children, rects):
            child.layout(rect)

    def _apply_split(self, lengths: List[Size], direction: Literal["vertical", "horizontal", "v", "h"]) -> List[Area]:
        if direction.lower() in ['vertical','v']:
            direction = 'vertical'
        elif direction.lower() in ['horizontal','h']:
            direction = 'horizontal'
        else:
            raise LayoutError(f"Split must be 'horizontal' or 'vertical', not: {direction}")

        if not lengths:
            raise SpecMismatchError(f"Must provide at least one {direction} length")

        self.children = []
        for l in lengths:
            self.add_child(Area(l))
        self.split_direction = direction
        return self.children

    def split_vertical(self, lengths: List[Size]) -> List[Area]:
        """Split this box into columns (children) with given widths."""
        return self._apply_split(lengths, "vertical")

    def split_horizontal(self, lengths: List[Size]) -> List[Area]:
        """Split this box into rows (children) with given heights."""
        return self._apply_split(lengths, "horizontal")

    def debug_rect(self, text: str|None=None):
        self.root.draw_rect(self.rect)
        if text:
           self.root.text(self.rect, f"{self.pos} - {text}")
        else:
            self.root.text(self.rect, f"{self.pos}")

    def add_table(self, df: pl.DataFrame) -> Table:
        return self.root.table(self.rect, df)

# Root container ----------------------------------------------------
class SlideRoot(Area):
    def __init__(self, prs: Presentation) -> None:
        super().__init__()
        self.prs = prs
        # TODO better logic to find/create empty slide by default
        self._slide = prs.slides.add_slide(prs.slide_layouts[6])

        if not isinstance(prs.slide_width, int):
            raise PPTXError(f"Unknown slide_width: {prs.slide_width}")
        if not isinstance(prs.slide_height, int):
            raise PPTXError(f"Unknown slide_height: {prs.slide_height}")

        self.slide_width = Emu(prs.slide_width)
        self.slide_height = Emu(prs.slide_height)

    def resolve(self) -> None:
        super().layout(Rect(Emu(0), Emu(0), self.slide_width, self.slide_height))

    def draw_rect(self,
        rect: Rect,
        *,
        fill: Optional[str] = None,
        line: str = "FF0000",
        line_width: int | Pt = 1
    ) -> None:
        shape = self._slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, rect.x, rect.y, rect.width, rect.height
        )
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill)
        shape.line.color.rgb = RGBColor.from_string(line)
        shape.line.width = line_width if isinstance(line_width, Pt) else Pt(line_width)


    def text(self, rect: Rect, text: str) -> None:
        shape = self._slide.shapes.add_textbox( rect.x, rect.y, rect.width, rect.height)
        shape.text = text

    def table(self, rect: Rect, df: pl.DataFrame) -> Table:
        rows = df.height+1
        cols = df.width
        gframe = self._slide.shapes.add_table(rows, cols, rect.x, rect.y, rect.width, rect.height)
        return Table(df, gframe)
