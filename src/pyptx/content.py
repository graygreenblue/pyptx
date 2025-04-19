import polars as pl
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import _Cell
from pptx.util import Pt

class Table:
    def __init__(self, df: pl.DataFrame, frame: GraphicFrame) -> None:
        self.df = df
        self._frame = frame
        self._table = frame.table

        self._write_header()
        self._write_data()
        for cell in self._table.iter_cells():
            self._format_cell(cell)

        self._auto_fontsize()


    def _write_header(self):
        for ci, hdr in enumerate(self.df.columns):
            self._table.cell(0,ci).text = str(hdr)

    def _write_data(self):
        for ri, row in enumerate(self.df.iter_rows(named = False), 1):
            for ci, val in enumerate(row):
                self._table.cell(ri, ci).text = str(val)


    def _format_cell(self,cell: _Cell):
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    def _auto_fontsize(self, max_size: float=40, min_size:float=1, step:float = 1):
        font_size = max_size
        prev_size = None

        while font_size >= min_size:
            for cell in self._table.iter_cells():
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(font_size)

            # Refresh layout
            current_height = self._frame.height
            if prev_size == current_height:
                break
            prev_size = current_height
            font_size -= step
