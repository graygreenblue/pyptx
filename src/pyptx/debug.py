"""
Debug helpers for **pptx_layout**.

Provide utilities to draw rectangles that match layout geometry on a
`python-pptx` slide so you can visually verify placement.
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

from .core import Rect, LayoutItem

# ------------------------------------------------------------------ #
# Drawing primitives
# ------------------------------------------------------------------ #
def draw_rect(slide, rect: Rect, *, fill: str | None = None,
              line: str = "FF0000", line_width: Pt = Pt(1)):
    """
    Add a rectangle shape to *slide* corresponding to *rect* (in EMUs).

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide.
    rect : Rect
        Rectangle in EMUs (x, y, width, height).
    fill : str | None, default None
        Hex RGB fill.  ``None`` → no fill.
    line : str, default "FF0000"
        Hex RGB outline colour.
    line_width : pptx.util.Pt, default 1 pt
        Outline thickness.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        rect.x, rect.y, rect.width, rect.height
    )

    # Fill
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill)

    # Outline
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = line_width
    return shape


def draw_layout(slide, root: LayoutItem, *,
                fill: str | None = None, line: str = "FF0000", line_width: Pt = Pt(1)):
    """
    Draw rectangles for *every* node in a resolved layout tree.

    Call *after* you've run ``root.resolve()``.
    """
    for node in root.walk():
        if node.rect is not None:
            draw_rect(slide, node.rect, fill=fill, line=line, line_width=line_width)


# ------------------------------------------------------------------ #
# Quick demo
# ------------------------------------------------------------------ #
def test_draw_rects(filename: str = "debug_rects_demo.pptx"):
    """
    Generate a PowerPoint file with a few rectangles for manual inspection.

    The demo is intentionally minimal; customise as needed.
    """
    from pptx import Presentation

    from .units import Ratio, Weight
    from .core import SlideLayout, Row, Column, Box

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Build a simple layout: ⅓ top band; remaining space split into 2 columns.
    root = SlideLayout(prs.slide_width, prs.slide_height)
    col = root.add(Column([Ratio(0.33), Weight(1)]))

    # Top band – single box
    col_top = col.add_box(Ratio(0.33))

    # Lower area – horizontal split
    lower_row = Row([Ratio(0.4), Weight(1)])
    col.add(lower_row)
    lower_row.add_box(Ratio(0.4))
    lower_row.add_box(Weight(1))

    root.resolve()
    draw_layout(slide, root, fill=None, line="00AAFF")

    prs.save(filename)
    return filename
